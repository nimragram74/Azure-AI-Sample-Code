#!/usr/bin/env python3
"""Generate the "Anthropic COE" presentation (PowerPoint) for Wipro.

Content mirrors the COE portal (anthropic-coe). Branding uses the portal's
Wipro corporate-blue theme so the deck and the app feel like one product.
Run:  python3 anthropic-coe/presentation/build_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ---- Brand palette (matches anthropic-coe/public/config/site.xml) ----------
PRIMARY      = RGBColor(0x1A, 0x56, 0xDB)   # corporate blue
PRIMARY_DARK = RGBColor(0x12, 0x41, 0xA8)
INK          = RGBColor(0x0C, 0x1B, 0x33)   # deep navy
CREAM        = RGBColor(0xEA, 0xF1, 0xFC)   # soft blue band
PAPER        = RGBColor(0xF8, 0xFA, 0xFE)
ACCENT       = RGBColor(0x00, 0xA3, 0x9A)   # teal
ACCENT_SOFT  = RGBColor(0xD6, 0xF2, 0xEF)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
MUTED        = RGBColor(0x5B, 0x66, 0x77)
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---- low-level helpers -----------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    return tb, tf


def para(tf, text, size, color, bold=False, first=False, align=PP_ALIGN.LEFT,
         space_after=6, italic=False, font=FONT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return p


def chip(s, x, y, text, fill=ACCENT_SOFT, fg=PRIMARY_DARK, w=None):
    w = w or Inches(2.2)
    sp = rect(s, x, y, w, Inches(0.34), fill, MSO_SHAPE.ROUNDED_RECTANGLE)
    sp.adjustments[0] = 0.5
    tf = sp.text_frame
    tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(10.5); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = fg
    return sp


def eyebrow(s, text, x=Inches(0.9), y=Inches(0.55), color=PRIMARY):
    _, tf = textbox(s, x, y, Inches(11), Inches(0.4))
    para(tf, text.upper(), 12, color, bold=True, first=True)


def title(s, text, x=Inches(0.9), y=Inches(0.95), w=Inches(11.5), color=INK, size=32):
    _, tf = textbox(s, x, y, w, Inches(1.0))
    para(tf, text, size, color, bold=True, first=True)


def footer(s, page):
    line = rect(s, Inches(0.9), Inches(7.02), Inches(11.53), Pt(1.2), CREAM)
    _, tf = textbox(s, Inches(0.9), Inches(7.05), Inches(9), Inches(0.35))
    para(tf, "Wipro  ·  Anthropic Center of Excellence", 9, MUTED, first=True)
    _, tf2 = textbox(s, Inches(11.4), Inches(7.05), Inches(1.0), Inches(0.35))
    para(tf2, str(page), 9, MUTED, first=True, align=PP_ALIGN.RIGHT)


def card(s, x, y, w, h, title_txt, body_txt, tag=None, accent=PRIMARY):
    c = rect(s, x, y, w, h, WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
    c.adjustments[0] = 0.06
    c.line.color.rgb = CREAM; c.line.width = Pt(1)
    bar = rect(s, x, y, Inches(0.09), h, accent, MSO_SHAPE.ROUNDED_RECTANGLE)
    bar.adjustments[0] = 0.5
    _, tf = textbox(s, x + Inches(0.28), y + Inches(0.16), w - Inches(0.5), h - Inches(0.3))
    para(tf, title_txt, 15, INK, bold=True, first=True, space_after=4)
    if tag:
        para(tf, tag, 9.5, accent, bold=True, space_after=5)
    para(tf, body_txt, 11, MUTED, space_after=0)
    return c


# ===========================================================================
# 1 — TITLE
# ===========================================================================
s = slide(); bg(s, INK)
rect(s, 0, 0, SW, Inches(0.16), PRIMARY)
rect(s, 0, SH - Inches(0.16), SW, Inches(0.16), ACCENT)
# logo mark
mark = rect(s, Inches(0.9), Inches(0.85), Inches(0.7), Inches(0.7), PRIMARY, MSO_SHAPE.ROUNDED_RECTANGLE)
mark.adjustments[0] = 0.28
mtf = mark.text_frame; mtf.paragraphs[0].alignment = PP_ALIGN.CENTER
mr = mtf.paragraphs[0].add_run(); mr.text = "A"; mr.font.size = Pt(28); mr.font.bold = True
mr.font.color.rgb = WHITE; mr.font.name = FONT
_, tf = textbox(s, Inches(1.75), Inches(0.9), Inches(8), Inches(0.7))
para(tf, "Anthropic COE", 20, WHITE, bold=True, first=True, space_after=0)
para(tf, "WIPRO · CENTER OF EXCELLENCE", 10, RGBColor(0x9F,0xB6,0xE8), bold=True)

_, tf = textbox(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(2.4))
para(tf, "Building the future of", 40, WHITE, bold=True, first=True, space_after=2)
para(tf, "enterprise AI with Claude", 40, RGBColor(0x6F,0xA8,0xFF), bold=True, space_after=14)
para(tf, "A single home for everything Anthropic at Wipro — training, certifications, "
        "engineering best practices, the latest updates, and proven blueprints to help "
        "our customers ship safe, useful AI.", 15, RGBColor(0xC7,0xD3,0xE8), space_after=0)

for i, (num, lbl) in enumerate([("4", "Learning tracks"), ("6", "Certifications"),
                                  ("40+", "Best practices"), ("12+", "Accelerators")]):
    x = Inches(0.9 + i * 2.55)
    box = rect(s, x, Inches(5.55), Inches(2.3), Inches(1.05), RGBColor(0x16,0x29,0x47), MSO_SHAPE.ROUNDED_RECTANGLE)
    box.adjustments[0] = 0.12
    _, tf = textbox(s, x, Inches(5.62), Inches(2.3), Inches(0.95), MSO_ANCHOR.MIDDLE)
    para(tf, num, 26, RGBColor(0x6F,0xA8,0xFF), bold=True, first=True, align=PP_ALIGN.CENTER, space_after=0)
    para(tf, lbl, 11, RGBColor(0xC7,0xD3,0xE8), align=PP_ALIGN.CENTER, space_after=0)


# ===========================================================================
# 2 — AGENDA
# ===========================================================================
s = slide(); bg(s, PAPER)
eyebrow(s, "Agenda"); title(s, "What we'll cover")
items = [
    ("01", "What is the Anthropic COE", "Our mission and what the portal brings together"),
    ("02", "Capabilities", "Everything the COE offers under one roof"),
    ("03", "Training & Certifications", "Role-based learning paths and stackable credentials"),
    ("04", "Updates & Best Practices", "Staying current and building things that work"),
    ("05", "How Wipro helps customers", "Accelerators and our engagement model"),
    ("06", "Get started", "How to engage the COE"),
]
for i, (n, t, d) in enumerate(items):
    col = i % 2; row = i // 2
    x = Inches(0.9 + col * 5.95); y = Inches(1.95 + row * 1.45)
    nb = rect(s, x, y, Inches(0.85), Inches(0.85), ACCENT_SOFT, MSO_SHAPE.ROUNDED_RECTANGLE)
    nb.adjustments[0] = 0.25
    ntf = nb.text_frame; ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
    np = ntf.paragraphs[0]; np.alignment = PP_ALIGN.CENTER
    nr = np.add_run(); nr.text = n; nr.font.size = Pt(18); nr.font.bold = True
    nr.font.color.rgb = PRIMARY; nr.font.name = FONT
    _, tf = textbox(s, x + Inches(1.05), y, Inches(4.6), Inches(0.95), MSO_ANCHOR.MIDDLE)
    para(tf, t, 15, INK, bold=True, first=True, space_after=2)
    para(tf, d, 10.5, MUTED, space_after=0)
footer(s, 2)


# ===========================================================================
# 3 — WHAT IS THE COE
# ===========================================================================
s = slide(); bg(s, PAPER)
rect(s, 0, 0, Inches(4.6), SH, INK)
_, tf = textbox(s, Inches(0.5), Inches(2.4), Inches(3.7), Inches(2.5))
para(tf, "What is the", 24, RGBColor(0x9F,0xB6,0xE8), bold=True, first=True, space_after=0)
para(tf, "Anthropic COE?", 24, WHITE, bold=True, space_after=0)
eyebrow(s, "Our mission", x=Inches(5.1))
_, tf = textbox(s, Inches(5.1), Inches(1.4), Inches(7.4), Inches(1.4))
para(tf, "Wipro's specialised practice for designing, building and scaling enterprise "
        "solutions on Claude — responsibly and at speed.", 18, INK, bold=True, first=True)
points = [
    ("Curated enablement", "Role-based learning and certifications that build real, applied skill."),
    ("Engineering excellence", "A living playbook of best practices for prompting, agents, evals and safety."),
    ("Customer impact", "Reusable accelerators that turn Claude's capabilities into measurable value."),
    ("Responsible by design", "Governance and guardrails baked into everything we build."),
]
for i, (t, d) in enumerate(points):
    y = Inches(2.95 + i * 0.97)
    dot = rect(s, Inches(5.1), y + Inches(0.05), Inches(0.22), Inches(0.22), ACCENT, MSO_SHAPE.OVAL)
    _, tf = textbox(s, Inches(5.5), y, Inches(7.0), Inches(0.9))
    para(tf, t, 13.5, PRIMARY_DARK, bold=True, first=True, space_after=1)
    para(tf, d, 11, MUTED, space_after=0)
footer(s, 3)


# ===========================================================================
# 4 — CAPABILITIES (6 cards)
# ===========================================================================
s = slide(); bg(s, PAPER)
eyebrow(s, "Capabilities"); title(s, "One portal. Every capability.")
caps = [
    ("Guided Learning", "Role-based tracks from first prompt to production agents."),
    ("Certifications", "Validate Claude & agent-building skills with recognised credentials."),
    ("Best Practices", "Battle-tested patterns for prompting, tools, evals, safety and cost."),
    ("Live Updates & Feeds", "Model releases, research and COE news, aggregated."),
    ("Customer Accelerators", "Blueprints that turn Claude into measurable customer value."),
    ("Responsible AI", "Governance, safety and evaluation guardrails by default."),
]
cw, ch = Inches(3.78), Inches(1.95)
for i, (t, d) in enumerate(caps):
    col = i % 3; row = i // 3
    x = Inches(0.9 + col * 3.95); y = Inches(2.0 + row * 2.25)
    card(s, x, y, cw, ch, t, d, accent=(PRIMARY if i % 2 == 0 else ACCENT))
footer(s, 4)


# ===========================================================================
# 5 — TRAINING (4 tracks)
# ===========================================================================
s = slide(); bg(s, CREAM)
eyebrow(s, "Training"); title(s, "Learning paths for every role")
tracks = [
    ("Claude Foundations", "FOUNDATIONAL · ~4 HRS", "Prompting, API fundamentals and responsible use. Intro to Claude, prompt engineering, model selection."),
    ("Building with Claude", "INTERMEDIATE · ~8 HRS", "Hands-on: tool use, RAG, structured outputs, streaming, prompt caching and cost optimisation."),
    ("Agents & the Agent SDK", "ADVANCED · ~10 HRS", "Multi-step agents, Model Context Protocol (MCP), the Claude Agent SDK and Claude Code."),
    ("Enterprise & Responsible AI", "LEADERSHIP · ~5 HRS", "Evals, safety, red-teaming, deployment patterns and the economics of AI adoption."),
]
cw, ch = Inches(5.78), Inches(1.95)
for i, (t, tag, d) in enumerate(tracks):
    col = i % 2; row = i // 2
    x = Inches(0.9 + col * 5.95); y = Inches(2.0 + row * 2.25)
    card(s, x, y, cw, ch, t, d, tag=tag, accent=(PRIMARY if i % 2 == 0 else ACCENT))
footer(s, 5)


# ===========================================================================
# 6 — CERTIFICATIONS (4 tiers)
# ===========================================================================
s = slide(); bg(s, PAPER)
eyebrow(s, "Certifications"); title(s, "Certify your Claude expertise")
_, tf = textbox(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.5))
para(tf, "Stackable credentials that prove real, applied skill — recognised across the practice.", 13, MUTED, first=True)
certs = [
    ("Certified Claude Associate", "ASSOCIATE · 2 WEEKS", "Prompting, API fundamentals, responsible use."),
    ("Certified Claude Developer", "DEVELOPER · 4 WEEKS", "Tool use, RAG, structured outputs, evals."),
    ("Certified Agent Engineer", "EXPERT · 6 WEEKS", "Multi-step agents, MCP, Claude Agent SDK."),
    ("Certified Solution Architect", "ARCHITECT · 8 WEEKS", "Enterprise architecture, security, scale."),
]
cw, ch = Inches(2.78), Inches(3.4)
for i, (t, tag, d) in enumerate(certs):
    x = Inches(0.9 + i * 2.95); y = Inches(2.45)
    card(s, x, y, cw, ch, t, d, tag=tag, accent=(PRIMARY if i % 2 == 0 else ACCENT))
    # tier number badge
    b = rect(s, x + Inches(0.28), y + ch - Inches(0.7), Inches(0.55), Inches(0.45), ACCENT_SOFT, MSO_SHAPE.ROUNDED_RECTANGLE)
    b.adjustments[0] = 0.3
    btf = b.text_frame; btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run(); br.text = f"L{i+1}"; br.font.size = Pt(13); br.font.bold = True
    br.font.color.rgb = PRIMARY_DARK; br.font.name = FONT
footer(s, 6)


# ===========================================================================
# 7 — UPDATES & BEST PRACTICES (split)
# ===========================================================================
s = slide(); bg(s, PAPER)
eyebrow(s, "Stay current · Build well"); title(s, "Latest updates & best practices")
# left: updates
lx = Inches(0.9)
hdr = rect(s, lx, Inches(1.9), Inches(5.6), Inches(0.5), INK, MSO_SHAPE.ROUNDED_RECTANGLE)
hdr.adjustments[0] = 0.25
htf = hdr.text_frame; htf.vertical_anchor = MSO_ANCHOR.MIDDLE
hp = htf.paragraphs[0]; hp.alignment = PP_ALIGN.CENTER
hr = hp.add_run(); hr.text = "LATEST UPDATES & FEEDS"; hr.font.size = Pt(12); hr.font.bold = True
hr.font.color.rgb = WHITE; hr.font.name = FONT
ups = [
    ("Claude Opus 4.8 & the 4.X family", "Most capable models — default for new builds."),
    ("Claude Agent SDK & MCP momentum", "Building agents is easier than ever."),
    ("Interpretability & alignment research", "Foundational to how we design safe systems."),
    ("Responsible Scaling Policy update", "How safety scales with capability."),
]
for i, (t, d) in enumerate(ups):
    y = Inches(2.6 + i * 1.0)
    rect(s, lx, y + Inches(0.04), Inches(0.09), Inches(0.7), ACCENT, MSO_SHAPE.ROUNDED_RECTANGLE).adjustments[0] = 0.5
    _, tf = textbox(s, lx + Inches(0.28), y, Inches(5.3), Inches(0.95))
    para(tf, t, 12.5, INK, bold=True, first=True, space_after=1)
    para(tf, d, 10.5, MUTED, space_after=0)
# right: best practices
rxp = Inches(6.95)
hdr2 = rect(s, rxp, Inches(1.9), Inches(5.45), Inches(0.5), PRIMARY, MSO_SHAPE.ROUNDED_RECTANGLE)
hdr2.adjustments[0] = 0.25
h2tf = hdr2.text_frame; h2tf.vertical_anchor = MSO_ANCHOR.MIDDLE
h2p = h2tf.paragraphs[0]; h2p.alignment = PP_ALIGN.CENTER
h2r = h2p.add_run(); h2r.text = "BEST PRACTICES FOR CLAUDE"; h2r.font.size = Pt(12); h2r.font.bold = True
h2r.font.color.rgb = WHITE; h2r.font.name = FONT
bps = [
    ("Be clear, direct & contextual", "State goal, audience and format."),
    ("Structure prompts with XML tags", "Reliable parsing of intent."),
    ("Start simple, add agency only when needed", "Workflow before full autonomy."),
    ("Define success & automate evals", "Measure quality, don't guess."),
    ("Cache context, right-size the model", "Cut latency and cost."),
]
for i, (t, d) in enumerate(bps):
    y = Inches(2.58 + i * 0.82)
    chk = rect(s, rxp, y + Inches(0.02), Inches(0.22), Inches(0.22), ACCENT, MSO_SHAPE.OVAL)
    _, tf = textbox(s, rxp + Inches(0.38), y, Inches(5.0), Inches(0.78))
    para(tf, t, 12, PRIMARY_DARK, bold=True, first=True, space_after=0)
    para(tf, d, 10, MUTED, space_after=0)
footer(s, 7)


# ===========================================================================
# 8 — HOW WIPRO HELPS CUSTOMERS (6 solutions)
# ===========================================================================
s = slide(); bg(s, PAPER)
eyebrow(s, "For customers"); title(s, "How Wipro helps customers with Claude")
sols = [
    ("AI Value Discovery", "Map high-value use cases to Claude's strengths.", "2–3 wks to a roadmap"),
    ("Agentic Automation", "Multi-step agents that automate knowledge work.", "Agent in 4–6 wks"),
    ("Enterprise Knowledge Assistant", "Secure RAG with citations and guardrails.", "Grounded day one"),
    ("Engineering Acceleration", "Claude Code & Agent SDK across the SDLC.", "Up to 40% faster"),
    ("Intelligent Customer Experience", "Support copilots and smart self-service.", "Higher CSAT"),
    ("Responsible AI & Governance", "Safety, evals, compliance and monitoring.", "Audit-ready"),
]
cw, ch = Inches(3.78), Inches(1.95)
for i, (t, d, out) in enumerate(sols):
    col = i % 3; row = i // 3
    x = Inches(0.9 + col * 3.95); y = Inches(2.0 + row * 2.25)
    card(s, x, y, cw, ch, t, d, accent=(PRIMARY if i % 2 == 0 else ACCENT))
    chip(s, x + Inches(0.28), y + ch - Inches(0.55), out, fill=ACCENT_SOFT, fg=PRIMARY_DARK, w=Inches(2.4))
footer(s, 8)


# ===========================================================================
# 9 — ENGAGEMENT MODEL (4 steps)
# ===========================================================================
s = slide(); bg(s, INK)
eyebrow(s, "Engagement model", color=RGBColor(0x9F,0xB6,0xE8))
title(s, "From idea to enterprise scale", color=WHITE)
steps = [
    ("1", "Discover", "Identify and prioritise high-value, feasible use cases."),
    ("2", "Prototype", "Build a Claude proof-of-value against real data in weeks."),
    ("3", "Industrialise", "Harden with evals, guardrails, security and CI/CD."),
    ("4", "Scale", "Roll out across the enterprise with COE governance."),
]
cw = Inches(2.78)
for i, (n, t, d) in enumerate(steps):
    x = Inches(0.9 + i * 2.95); y = Inches(2.4)
    box = rect(s, x, y, cw, Inches(3.1), RGBColor(0x16,0x29,0x47), MSO_SHAPE.ROUNDED_RECTANGLE)
    box.adjustments[0] = 0.06
    nb = rect(s, x + Inches(0.3), y + Inches(0.3), Inches(0.7), Inches(0.7),
              PRIMARY if i % 2 == 0 else ACCENT, MSO_SHAPE.ROUNDED_RECTANGLE)
    nb.adjustments[0] = 0.28
    ntf = nb.text_frame; ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
    npp = ntf.paragraphs[0]; npp.alignment = PP_ALIGN.CENTER
    nr = npp.add_run(); nr.text = n; nr.font.size = Pt(22); nr.font.bold = True
    nr.font.color.rgb = WHITE; nr.font.name = FONT
    _, tf = textbox(s, x + Inches(0.3), y + Inches(1.2), cw - Inches(0.6), Inches(1.8))
    para(tf, t, 17, WHITE, bold=True, first=True, space_after=4)
    para(tf, d, 11, RGBColor(0xC7,0xD3,0xE8), space_after=0)
    if i < 3:
        _, atf = textbox(s, x + cw - Inches(0.1), y + Inches(1.2), Inches(0.5), Inches(0.5))
        para(atf, "→", 22, RGBColor(0x6F,0xA8,0xFF), bold=True, first=True, align=PP_ALIGN.CENTER)


# ===========================================================================
# 10 — CALL TO ACTION
# ===========================================================================
s = slide(); bg(s, PAPER)
band = rect(s, Inches(0.9), Inches(1.5), Inches(11.53), Inches(4.5), PRIMARY, MSO_SHAPE.ROUNDED_RECTANGLE)
band.adjustments[0] = 0.05
rect(s, Inches(0.9), Inches(1.5), Inches(11.53), Inches(0.16), ACCENT, MSO_SHAPE.ROUNDED_RECTANGLE).adjustments[0] = 0.5
_, tf = textbox(s, Inches(1.6), Inches(2.2), Inches(10.1), Inches(3.2), MSO_ANCHOR.MIDDLE)
para(tf, "Ready to explore what Claude can do", 30, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER, space_after=2)
para(tf, "for your organisation?", 30, WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=16)
para(tf, "Talk to the Wipro Anthropic Center of Excellence", 16, RGBColor(0xDCE6,0xFF & 0xFF, 0xFF) if False else WHITE, align=PP_ALIGN.CENTER, space_after=14)
btn = rect(s, Inches(4.95), Inches(4.55), Inches(3.4), Inches(0.7), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
btn.adjustments[0] = 0.5
btf = btn.text_frame; btf.vertical_anchor = MSO_ANCHOR.MIDDLE
bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
br = bp.add_run(); br.text = "anthropic.coe@wipro.com"; br.font.size = Pt(14); br.font.bold = True
br.font.color.rgb = PRIMARY_DARK; br.font.name = FONT
_, tf = textbox(s, Inches(0.9), Inches(6.4), Inches(11.53), Inches(0.5))
para(tf, "© 2026 Wipro Limited.  Anthropic, Claude and related marks are property of Anthropic, PBC.",
     10, MUTED, first=True, align=PP_ALIGN.CENTER)


out = os.path.join(os.path.dirname(__file__), "Anthropic-COE-Portal.pptx")
prs.save(out)
print("Saved", out, "—", len(prs.slides._sldIdLst), "slides")
